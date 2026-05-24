from __future__ import annotations

from pathlib import Path

from privcage.errors import ParseError
from privcage.models import ParseResult


def parse_docx(path: Path) -> ParseResult:
    try:
        from docx import Document
    except ImportError as exc:
        raise ParseError("python-docx is required for docx files") from exc

    document = Document(str(path))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            lines.append(" | ".join(cell.text.strip() for cell in row.cells))
    return ParseResult(text="\n\n".join(lines), source_kind="docx")


def parse_xlsx(path: Path) -> ParseResult:
    try:
        import openpyxl
    except ImportError as exc:
        raise ParseError("openpyxl is required for xlsx files") from exc

    workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
    chunks: list[str] = []
    for sheet in workbook.worksheets:
        chunks.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                chunks.append(" | ".join(values))
    return ParseResult(text="\n".join(chunks), source_kind="xlsx")


def parse_pptx(path: Path) -> ParseResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParseError("python-pptx is required for pptx files") from exc

    presentation = Presentation(str(path))
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        chunks.append(f"## Slide {slide_index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                chunks.append(text.strip())
    return ParseResult(text="\n\n".join(chunks), source_kind="pptx")
